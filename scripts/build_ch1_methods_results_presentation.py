#!/usr/bin/env python3

from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Iterable

os.environ.setdefault(
    "MPLCONFIGDIR",
    str(
        Path(__file__).resolve().parents[1]
        / "reports"
        / "ch1_methods_results_assets"
        / ".mplconfig"
    ),
)

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = ROOT / "reports"
ASSETS_DIR = REPORTS_DIR / "ch1_methods_results_assets"
CLUSTER_DIR = ROOT / "cluster-results" / "chapter1_true_results"
OUTPUT_PPTX = REPORTS_DIR / "ch1_methods_results_presentation.pptx"

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
BG = RGBColor(255, 255, 255)
TEXT = RGBColor(30, 41, 59)
SUBTEXT = RGBColor(71, 85, 105)
ACCENT = RGBColor(15, 118, 110)
ACCENT_LIGHT = RGBColor(240, 253, 250)
ACCENT_SOFT = RGBColor(204, 251, 241)
RUST = RGBColor(194, 65, 12)
RUST_LIGHT = RGBColor(255, 247, 237)
PANEL = RGBColor(248, 250, 252)
LINE = RGBColor(226, 232, 240)
MUTED = RGBColor(100, 116, 139)


def ensure_dirs() -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def load_csv(rel_path: str) -> pd.DataFrame:
    return pd.read_csv(CLUSTER_DIR / rel_path)


def fmt_int(value: int | float) -> str:
    return f"{int(round(float(value))):,}"


def fmt_pct(value: float, digits: int = 1) -> str:
    return f"{100 * float(value):.{digits}f}%"


def add_run(title: str, paragraph, font_size: int = 16, bold: bool = False, color=TEXT):
    run = paragraph.add_run()
    run.text = title
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def style_text_frame(text_frame, font_size: int = 14, color=TEXT) -> None:
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(4)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(3)
    text_frame.margin_bottom = Pt(3)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT


def add_card(slide, left, top, width, height, title: str, lines: list[str], fill=PANEL, title_fill=None) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    title_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.05),
        top + Inches(0.05),
        width - Inches(0.10),
        Inches(0.36),
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = title_fill or ACCENT_LIGHT
    title_bar.line.color.rgb = title_fill or ACCENT_LIGHT
    tf = title_bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = BODY_FONT
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = TEXT

    text_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.48), width - Inches(0.24), height - Inches(0.58))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_after = Pt(3)
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(11)
            run.font.color.rgb = TEXT


def add_metric_card(slide, left, top, width, height, label: str, value: str, note: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT
    shape.line.color.rgb = ACCENT_SOFT
    shape.line.width = Pt(1)

    label_box = slide.shapes.add_textbox(left + Inches(0.10), top + Inches(0.08), width - Inches(0.20), Inches(0.24))
    tf = label_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = SUBTEXT

    value_box = slide.shapes.add_textbox(left + Inches(0.10), top + Inches(0.28), width - Inches(0.20), Inches(0.34))
    tf = value_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.name = TITLE_FONT
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    note_box = slide.shapes.add_textbox(left + Inches(0.10), top + Inches(0.67), width - Inches(0.20), height - Inches(0.74))
    tf = note_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = note
    run.font.name = BODY_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT


def add_interpretation_box(slide, text: str, caveat: bool = False) -> None:
    left = Inches(0.48)
    width = Inches(12.32)
    top = Inches(6.78)
    height = Inches(0.45)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RUST_LIGHT if caveat else ACCENT_LIGHT
    shape.line.color.rgb = RUST if caveat else ACCENT_SOFT
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    label = "Caveat: " if caveat else "Interpretation: "
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RUST if caveat else ACCENT
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.color.rgb = TEXT


def add_source_note(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.18), Inches(10.8), Inches(0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Source: {text}"
    run.font.name = BODY_FONT
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED


def add_slide_number(slide, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(12.0), Inches(7.14), Inches(0.8), Inches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = label
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_header(slide, title: str, subtitle: str | None = None, section: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.28), Inches(9.8), Inches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.54), Inches(0.70), Inches(10.8), Inches(0.24))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = BODY_FONT
        run.font.size = Pt(11)
        run.font.color.rgb = SUBTEXT

    if section:
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(0.30), Inches(1.25), Inches(0.34))
        chip.fill.solid()
        chip.fill.fore_color.rgb = ACCENT_LIGHT if section != "Appendix" else RUST_LIGHT
        chip.line.color.rgb = ACCENT_SOFT if section != "Appendix" else RGBColor(253, 186, 116)
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = section
        run.font.name = BODY_FONT
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = ACCENT if section != "Appendix" else RUST


def add_image_contain(slide, path: Path, left, top, width, height):
    img = Image.open(path)
    img_w, img_h = img.size
    box_w = width / 914400
    box_h = height / 914400
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio >= box_ratio:
        final_w = width
        final_h = int(width / img_ratio)
        final_left = left
        final_top = top + int((height - final_h) / 2)
    else:
        final_h = height
        final_w = int(height * img_ratio)
        final_top = top
        final_left = left + int((width - final_w) / 2)
    slide.shapes.add_picture(str(path), final_left, final_top, width=final_w, height=final_h)


def add_table(slide, left, top, width, height, df: pd.DataFrame, font_size: int = 10, header_fill=ACCENT_LIGHT):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    col_width = int(width / cols)
    for col_idx in range(cols):
        table.columns[col_idx].width = col_width
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = TEXT
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(df.iloc[row_idx, col_idx])
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = BODY_FONT
                run.font.size = Pt(font_size)
                run.font.color.rgb = TEXT
    return table


def export_label_diagnostics(cohort_summary: pd.DataFrame) -> Path:
    out_path = ASSETS_DIR / "slide04_label_diagnostics.png"
    horizon_rows = cohort_summary[cohort_summary["summary_group"] == "labels"].copy()
    horizon_rows["horizon_h"] = horizon_rows["horizon_h"].astype(int)
    plot_df = (
        horizon_rows.pivot_table(index="horizon_h", columns="metric", values="value", aggfunc="first")
        .reset_index()
        .sort_values("horizon_h")
    )
    plot_df["positive_prevalence"] = plot_df["positive_labels"] / plot_df["labelable_instances"]

    plt.style.use("seaborn-v0_8-whitegrid")
    fig, axes = plt.subplots(1, 2, figsize=(10.5, 3.8), dpi=180)

    x = plot_df["horizon_h"].astype(str)
    axes[0].bar(x, plot_df["labelable_instances"], color="#0f766e", label="Labelable")
    axes[0].bar(
        x,
        plot_df["unlabeled_instances"],
        bottom=plot_df["labelable_instances"],
        color="#cbd5e1",
        label="Unlabeled",
    )
    axes[0].set_title("Instance status by horizon", fontsize=11, weight="bold")
    axes[0].set_ylabel("Rows")
    axes[0].legend(frameon=False, fontsize=8, loc="upper right")
    axes[0].tick_params(labelsize=8)

    axes[1].plot(x, 100 * plot_df["positive_prevalence"], marker="o", color="#c2410c", linewidth=2)
    axes[1].set_title("Positive prevalence among labelable rows", fontsize=11, weight="bold")
    axes[1].set_ylabel("Percent")
    axes[1].tick_params(labelsize=8)
    for idx, val in enumerate(plot_df["positive_prevalence"]):
        axes[1].annotate(f"{100*val:.1f}%", (idx, 100 * val), textcoords="offset points", xytext=(0, 8), ha="center", fontsize=8)

    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def export_baseline_overview(combined_metrics: pd.DataFrame) -> Path:
    out_path = ASSETS_DIR / "slide07_baseline_overview.png"
    plot_df = combined_metrics[combined_metrics["split"] == "test"].copy()
    plot_df["horizon_h"] = plot_df["horizon_h"].astype(int)
    plot_df = plot_df.sort_values(["model_name", "horizon_h"])

    palette = {"logistic_regression": "#0f766e", "xgboost": "#1d4ed8"}
    labels = {"logistic_regression": "Logistic regression", "xgboost": "XGBoost"}

    plt.style.use("seaborn-v0_8-whitegrid")
    fig, axes = plt.subplots(1, 2, figsize=(10.5, 3.8), dpi=180)
    for model_name, sub_df in plot_df.groupby("model_name"):
        axes[0].plot(sub_df["horizon_h"], sub_df["auroc"], marker="o", linewidth=2, color=palette[model_name], label=labels[model_name])
        axes[1].plot(sub_df["horizon_h"], sub_df["calibration_slope"], marker="o", linewidth=2, color=palette[model_name], label=labels[model_name])
    axes[0].set_title("Test AUROC by horizon", fontsize=11, weight="bold")
    axes[0].set_xlabel("Horizon (h)")
    axes[0].set_ylim(0.74, 0.89)
    axes[0].tick_params(labelsize=8)

    axes[1].axhline(1.0, color="#64748b", linestyle="--", linewidth=1)
    axes[1].set_title("Calibration slope by horizon", fontsize=11, weight="bold")
    axes[1].set_xlabel("Horizon (h)")
    axes[1].set_ylim(0.90, 1.40)
    axes[1].tick_params(labelsize=8)
    axes[1].legend(frameon=False, fontsize=8, loc="upper right")

    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def export_hard_case_burden(hard_case_summary: pd.DataFrame) -> Path:
    out_path = ASSETS_DIR / "slide09_hard_case_burden.png"
    plot_df = hard_case_summary.copy()
    plot_df["horizon_h"] = plot_df["horizon_h"].astype(int)
    plot_df = plot_df.sort_values("horizon_h")

    plt.style.use("seaborn-v0_8-whitegrid")
    fig, ax = plt.subplots(figsize=(7.3, 3.8), dpi=180)
    bars = ax.bar(
        plot_df["horizon_h"].astype(str),
        100 * plot_df["pct_fatal_hard_cases"],
        color="#0f766e",
        width=0.62,
    )
    ax.set_title("Low-predicted fatal share by horizon", fontsize=12, weight="bold")
    ax.set_ylabel("Percent of fatal stays")
    ax.set_ylim(0, 28)
    ax.tick_params(labelsize=9)
    for bar, (_, row) in zip(bars, plot_df.iterrows()):
        ax.annotate(
            f"{100*row['pct_fatal_hard_cases']:.1f}%\n{int(row['n_hard_cases'])}/{int(row['n_fatal_last_points'])}",
            (bar.get_x() + bar.get_width() / 2, bar.get_height()),
            textcoords="offset points",
            xytext=(0, 6),
            ha="center",
            fontsize=8,
        )
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def load_data() -> dict:
    cohort_summary = load_csv("cohort/chapter1_cohort_summary.csv")
    site_eligibility = load_csv("cohort/chapter1_site_eligibility.csv")
    counts_by_hospital = load_csv("cohort/chapter1_counts_by_hospital.csv")
    model_ready = load_csv("model_ready/chapter1_primary_readiness_summary.csv")
    stay_split = load_csv("splits/chapter1_stay_split_summary.csv")
    primary_split = load_csv("splits/chapter1_primary_split_summary.csv")
    combined_metrics = load_csv("evaluation/asic/baselines/primary_medians/combined_metrics.csv")
    reporting_split = load_csv("evaluation/asic/baselines/primary_medians/reporting_split_summary.csv")
    site_summary = load_csv("evaluation/asic/baselines/primary_medians/combined_primary_site_summary.csv")
    hard_case_summary = load_csv(
        "evaluation/asic/hard_cases/primary_medians/logistic_regression/horizon_hard_case_summary.csv"
    )
    comparison_table = load_csv(
        "evaluation/asic/hard_cases/primary_medians/logistic_regression/asic_hard_case_comparison/comparison_table.csv"
    )
    agreement = load_csv(
        "evaluation/asic/hard_cases/primary_medians/agreement/logistic_regression_vs_xgboost/horizon_hard_case_agreement_summary.csv"
    )
    pairwise_overlap = load_csv("evaluation/asic/horizon_dependence/overlap/pairwise_overlap.csv")
    obs_qc = load_csv("observation_process/chapter1_observation_process_qc_summary.csv")
    agg_metrics = load_csv("temporal_preview/asic/aggregation_16h/comparison/aggregation_comparison_metrics.csv")
    with open(CLUSTER_DIR / "evaluation/asic/horizon_dependence/final/run_manifest.json", "r", encoding="utf-8") as f:
        horizon_manifest = json.load(f)
    return {
        "cohort_summary": cohort_summary,
        "site_eligibility": site_eligibility,
        "counts_by_hospital": counts_by_hospital,
        "model_ready": model_ready,
        "stay_split": stay_split,
        "primary_split": primary_split,
        "combined_metrics": combined_metrics,
        "reporting_split": reporting_split,
        "site_summary": site_summary,
        "hard_case_summary": hard_case_summary,
        "comparison_table": comparison_table,
        "agreement": agreement,
        "pairwise_overlap": pairwise_overlap,
        "obs_qc": obs_qc,
        "agg_metrics": agg_metrics,
        "horizon_manifest": horizon_manifest,
    }


def cohort_value(df: pd.DataFrame, summary_group: str, metric: str, horizon: int | None = None) -> float:
    subset = df[(df["summary_group"] == summary_group) & (df["metric"] == metric)]
    if horizon is not None:
        subset = subset[subset["horizon_h"] == horizon]
    return float(subset.iloc[0]["value"])


def build_presentation(data: dict, assets: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cohort_summary = data["cohort_summary"]
    site_eligibility = data["site_eligibility"]
    counts_by_hospital = data["counts_by_hospital"]
    model_ready = data["model_ready"]
    stay_split = data["stay_split"]
    primary_split = data["primary_split"]
    combined_metrics = data["combined_metrics"]
    site_summary = data["site_summary"]
    hard_case_summary = data["hard_case_summary"]
    comparison_table = data["comparison_table"]
    agreement = data["agreement"]
    pairwise_overlap = data["pairwise_overlap"]
    obs_qc = data["obs_qc"]
    agg_metrics = data["agg_metrics"]
    horizon_manifest = data["horizon_manifest"]

    input_hospitals = int(cohort_value(cohort_summary, "cohort", "input_hospitals"))
    retained_hospitals = int(cohort_value(cohort_summary, "cohort", "retained_hospitals"))
    input_stays = int(cohort_value(cohort_summary, "cohort", "input_stays"))
    retained_stays = int(cohort_value(cohort_summary, "cohort", "retained_stays"))
    valid_total = int(cohort_value(cohort_summary, "instances", "valid_prediction_instances_total"))
    valid_per_h = int(cohort_value(cohort_summary, "instances", "valid_prediction_instances", 24))
    labelable_24 = int(cohort_value(cohort_summary, "labels", "labelable_instances", 24))
    pos_24 = int(cohort_value(cohort_summary, "labels", "positive_labels", 24))
    neg_24 = int(cohort_value(cohort_summary, "labels", "negative_labels", 24))
    unlabeled_24 = int(cohort_value(cohort_summary, "labels", "unlabeled_instances", 24))

    stage_site = int(counts_by_hospital["after_site_level_exclusion"].sum())
    stage_mech = int(counts_by_hospital["after_mech_vent_ge_24h_qc_exclusion"].sum())
    retained_counts = (
        counts_by_hospital.loc[counts_by_hospital["final_retained_stays"] > 0, ["hospital_id", "final_retained_stays"]]
        .rename(columns={"hospital_id": "Hospital", "final_retained_stays": "Retained stays"})
    )
    retained_counts["Retained stays"] = retained_counts["Retained stays"].map(fmt_int)

    split_overall = stay_split[stay_split["summary_level"] == "overall"].copy()
    split_table = split_overall[["split", "stay_count", "positive_stays", "label_prevalence"]].copy()
    split_table.columns = ["Split", "Stays", "Mortality+", "Stay mortality"]
    split_table["Split"] = split_table["Split"].str.capitalize()
    split_table["Stays"] = split_table["Stays"].map(fmt_int)
    split_table["Mortality+"] = split_table["Mortality+"].map(fmt_int)
    split_table["Stay mortality"] = split_table["Stay mortality"].map(lambda x: fmt_pct(x, 2))

    split_24 = primary_split[
        (primary_split["summary_level"] == "overall_horizon")
        & (primary_split["split"].isin(["train", "validation", "test"]))
        & (primary_split["horizon_h"] == 24)
    ][["split", "instance_count", "positive_labels", "label_prevalence"]].copy()
    split_24.columns = ["Split", "24h rows", "24h positives", "24h event rate"]
    split_24["Split"] = split_24["Split"].str.capitalize()
    split_24["24h rows"] = split_24["24h rows"].map(fmt_int)
    split_24["24h positives"] = split_24["24h positives"].map(fmt_int)
    split_24["24h event rate"] = split_24["24h event_rate"] if "24h event_rate" in split_24.columns else split_24["24h event rate"]
    split_24["24h event rate"] = primary_split[
        (primary_split["summary_level"] == "overall_horizon")
        & (primary_split["split"].isin(["train", "validation", "test"]))
        & (primary_split["horizon_h"] == 24)
    ]["label_prevalence"].map(lambda x: fmt_pct(x, 2)).tolist()

    metrics_24 = combined_metrics[
        (combined_metrics["split"] == "test") & (combined_metrics["horizon_h"] == 24)
    ][["model_name", "sample_count", "event_count", "auroc", "auprc", "calibration_slope", "brier_score"]].copy()
    metrics_24["Model"] = metrics_24["model_name"].map(
        {"logistic_regression": "Logistic regression", "xgboost": "XGBoost"}
    )
    metrics_24["Rows"] = metrics_24["sample_count"].map(fmt_int)
    metrics_24["Events"] = metrics_24["event_count"].map(fmt_int)
    metrics_24["AUROC"] = metrics_24["auroc"].map(lambda x: f"{x:.3f}")
    metrics_24["AUPRC"] = metrics_24["auprc"].map(lambda x: f"{x:.3f}")
    metrics_24["Slope"] = metrics_24["calibration_slope"].map(lambda x: f"{x:.3f}")
    metrics_24["Brier"] = metrics_24["brier_score"].map(lambda x: f"{x:.4f}")
    metrics_24_table = metrics_24[["Model", "Rows", "Events", "AUROC", "AUPRC", "Slope", "Brier"]]

    thresholds = hard_case_summary[["horizon_h", "nonfatal_q75_threshold"]].copy()
    thresholds["Horizon"] = thresholds["horizon_h"].astype(int).astype(str) + "h"
    thresholds["Q75 threshold"] = thresholds["nonfatal_q75_threshold"].map(lambda x: f"{x:.4f}")
    thresholds = thresholds[["Horizon", "Q75 threshold"]]

    site_24_logistic = site_summary[
        (site_summary["model_name"] == "logistic_regression")
        & (site_summary["split"] == "test")
        & (site_summary["horizon_h"] == 24)
    ][["hospital_id", "event_count", "auroc", "auprc", "calibration_slope"]].copy()
    site_24_logistic.columns = ["Hospital", "Events", "AUROC", "AUPRC", "Slope"]
    site_24_logistic["Events"] = site_24_logistic["Events"].map(fmt_int)
    for col in ["AUROC", "AUPRC", "Slope"]:
        site_24_logistic[col] = site_24_logistic[col].map(lambda x: f"{x:.3f}")

    full_metrics = combined_metrics[combined_metrics["split"] == "test"].copy()
    full_metrics["Model"] = full_metrics["model_name"].map(
        {"logistic_regression": "Logistic", "xgboost": "XGBoost"}
    )
    full_metrics["Horizon"] = full_metrics["horizon_h"].astype(int).astype(str) + "h"
    full_metrics["Rows"] = full_metrics["sample_count"].map(fmt_int)
    full_metrics["Events"] = full_metrics["event_count"].map(fmt_int)
    full_metrics["AUROC"] = full_metrics["auroc"].map(lambda x: f"{x:.3f}")
    full_metrics["AUPRC"] = full_metrics["auprc"].map(lambda x: f"{x:.3f}")
    full_metrics["Slope"] = full_metrics["calibration_slope"].map(lambda x: f"{x:.3f}")
    full_metrics["Brier"] = full_metrics["brier_score"].map(lambda x: f"{x:.4f}")
    full_metrics = full_metrics[["Model", "Horizon", "Rows", "Events", "AUROC", "AUPRC", "Slope", "Brier"]]
    full_metrics_log = full_metrics[full_metrics["Model"] == "Logistic"].reset_index(drop=True)
    full_metrics_xgb = full_metrics[full_metrics["Model"] == "XGBoost"].reset_index(drop=True)

    agreement_table = agreement[
        ["horizon_h", "n_logistic_hard", "n_xgb_hard", "n_both_hard", "jaccard_hard_case_overlap"]
    ].copy()
    agreement_table.columns = ["Horizon", "Logistic hard", "XGBoost hard", "Overlap", "Jaccard"]
    agreement_table["Horizon"] = agreement_table["Horizon"].astype(int).astype(str) + "h"
    agreement_table["Logistic hard"] = agreement_table["Logistic hard"].map(fmt_int)
    agreement_table["XGBoost hard"] = agreement_table["XGBoost hard"].map(fmt_int)
    agreement_table["Overlap"] = agreement_table["Overlap"].map(fmt_int)
    agreement_table["Jaccard"] = agreement_table["Jaccard"].map(lambda x: f"{x:.3f}")

    variable_audit_table = pd.DataFrame(
        [
            ["Age", "Not ready", "0/1,682 exact age", "Only age_group exists"],
            ["Sex", "Ready", "1,682/1,682", "Static join complete"],
            ["Disease group", "Ready", "1,682/1,682", "Hierarchy-sensitive ICD bag"],
            ["Prediction time", "Ready", "1,682/1,682", "Hours from ICU admission proxy"],
            ["Respiratory", "Ready", "1,233/1,682 PF", "Current-block PF when present"],
            ["Hemodynamic", "Ready with fallback", "1,642/1,682 MAP", "MAP cleaner than vasopressor mapping"],
            ["Renal", "Ready", "1,650/1,682 creatinine", "Major LOCF dependence"],
            ["Ventilation", "Ready", "1,449/1,682 PEEP", "Mostly current-block"],
        ],
        columns=["Family", "Status", "Completeness", "Note"],
    )

    sofa_points = [
        "Respiratory scorable in 73%; coagulation 92% but largely LOCF-based.",
        "Liver is 71% available, with 51% depending on 48h LOCF.",
        "MAP is 98% available, but standard vasopressor inputs are absent.",
        "No GCS field and no urine output field were found.",
        "Complete-case coverage across partially represented organs is 49% even after LOCF.",
    ]

    obs_summary_table = pd.DataFrame(
        [
            ["Usable 8h blocks", fmt_int(obs_qc.loc[obs_qc["metric"] == "block_rows_total", "value"].iloc[0]), ""],
            ["Core groups observed = 4", fmt_pct(obs_qc.loc[obs_qc["metric"] == "value_4_proportion", "value"].iloc[0], 1), ""],
            ["Core groups observed = 3", fmt_pct(obs_qc.loc[obs_qc["metric"] == "value_3_proportion", "value"].iloc[0], 1), ""],
            ["Median TSL HR", "0.25 h", "near-real-time"],
            ["Median TSL BP", "0.25 h", "near-real-time"],
            ["Median TSL Resp", "0.50 h", "less dense than HR/BP/Oxy"],
            ["Median TSL Oxy", "0.25 h", "near-real-time"],
        ],
        columns=["Metric", "Value", "Comment"],
    )

    preview_24 = agg_metrics[
        (agg_metrics["horizon_h"] == 24) & (agg_metrics["selected_split"] == "test")
    ].copy()
    preview_pivot = preview_24.pivot_table(index="model_name", columns="aggregation", values=["auroc", "auprc", "calibration_slope", "brier_score"])
    temporal_table = []
    for model_name, label in [("logistic_regression", "Logistic"), ("xgboost", "XGBoost")]:
        temporal_table.append(
            [
                label,
                f"{preview_pivot.loc[model_name, ('auroc', '8h')]:.3f} -> {preview_pivot.loc[model_name, ('auroc', '16h')]:.3f}",
                f"{preview_pivot.loc[model_name, ('auprc', '8h')]:.3f} -> {preview_pivot.loc[model_name, ('auprc', '16h')]:.3f}",
                f"{preview_pivot.loc[model_name, ('calibration_slope', '8h')]:.3f} -> {preview_pivot.loc[model_name, ('calibration_slope', '16h')]:.3f}",
                f"{preview_pivot.loc[model_name, ('brier_score', '8h')]:.3f} -> {preview_pivot.loc[model_name, ('brier_score', '16h')]:.3f}",
            ]
        )
    temporal_table = pd.DataFrame(
        temporal_table,
        columns=["Model", "AUROC", "AUPRC", "Slope", "Brier"],
    )

    overlap_mean = pairwise_overlap["jaccard_index"].mean()
    overlap_24_48 = float(pairwise_overlap[(pairwise_overlap["horizon_a"] == "24h") & (pairwise_overlap["horizon_b"] == "48h")]["jaccard_index"].iloc[0])
    overlap_24_72 = float(pairwise_overlap[(pairwise_overlap["horizon_a"] == "24h") & (pairwise_overlap["horizon_b"] == "72h")]["jaccard_index"].iloc[0])

    map_row = comparison_table[comparison_table["variable"] == "map_last"].iloc[0]
    pf_row = comparison_table[comparison_table["variable"] == "pf_ratio_last"].iloc[0]
    peep_row = comparison_table[comparison_table["variable"] == "peep_last"].iloc[0]
    pred_time_row = comparison_table[comparison_table["variable"] == "prediction_time_h"].iloc[0]
    disease_row = comparison_table[comparison_table["variable"] == "disease_group"].iloc[0]
    hospital_row = comparison_table[comparison_table["variable"] == "hospital_id"].iloc[0]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Chapter 1 Methods And Results", "Current ASIC analysis: stand-alone methods/results presentation for the true cluster-result bundle", "Main")
    add_textbox(
        slide,
        Inches(0.68),
        Inches(1.30),
        Inches(7.2),
        Inches(0.6),
        "Goal",
        font_size=22,
        bold=True,
    )
    add_bullets(
        slide,
        Inches(0.72),
        Inches(1.84),
        Inches(7.0),
        Inches(2.6),
        [
            "Review the frozen Chapter 1 design choices and exact operational implementation.",
            "Show how hospitals, stays, rows, and labels are retained or dropped.",
            "Summarize the current ASIC baseline, hard-case, and horizon results from cluster-results.",
            "Keep all interpretation measurement-bound to the recorded feature set, charting process, and temporal aggregation.",
        ],
        font_size=16,
    )
    add_card(
        slide,
        Inches(8.35),
        Inches(1.58),
        Inches(4.15),
        Inches(2.2),
        "Talk boundaries",
        [
            "No PhD framing, proposal alignment, or thesis storyline.",
            "No biological subtype claims, causal claims, or irreducible-randomness claims.",
            "Empirical source priority: cluster-results over older local artifacts.",
        ],
    )
    add_metric_card(slide, Inches(8.35), Inches(4.10), Inches(1.95), Inches(1.28), "Retained stays", fmt_int(retained_stays), "Current ASIC retained cohort")
    add_metric_card(slide, Inches(10.55), Inches(4.10), Inches(1.95), Inches(1.28), "24h test rows", fmt_int(33676), "Binary-evaluable test split")
    add_interpretation_box(slide, "This is a bounded review of near-term in-ICU mortality risk structure under the current ASIC representation, not a generic mortality-prediction talk.")
    add_source_note(slide, "ch1_methods_results_deck_overview.md; chapter1_analysis_spec_frozen_v1.md")
    add_slide_number(slide, "1")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Frozen Study Setup", "Frozen v1 analysis contract before ASIC result interpretation", "Methods")
    card_w = Inches(3.95)
    card_h = Inches(1.32)
    xs = [Inches(0.58), Inches(4.68), Inches(8.78)]
    ys = [Inches(1.35), Inches(3.02)]
    add_card(slide, xs[0], ys[0], card_w, card_h, "Data", ["Development dataset: ASIC", "External-validation target: MIMIC-IV", "This deck shows ASIC only"])
    add_card(slide, xs[1], ys[0], card_w, card_h, "Endpoint", ["Primary outcome: in-ICU mortality", "Hard-case analysis is defined on fatal stays", "No event-timed death label is available"])
    add_card(slide, xs[2], ys[0], card_w, card_h, "Horizons", ["Primary: 24h", "Main contrast: 48h", "Sensitivities: 8h, 16h, 72h"])
    add_card(slide, xs[0], ys[1], card_w, card_h, "Unit", ["Patient-time prediction instances", "Completed 8h blocks", "Operational first-stay ICU analysis"])
    add_card(slide, xs[1], ys[1], card_w, card_h, "Baseline models", ["Logistic regression", "XGBoost", "Calibration-aware evaluation"])
    add_card(slide, xs[2], ys[1], card_w, card_h, "Explicit non-claims", ["No biological death classes", "No irreducible stochasticity claim", "No causal attribution"])
    add_interpretation_box(slide, "The core scientific choices were frozen early enough that later ASIC findings can be read against a stable design contract rather than a post hoc analysis setup.")
    add_source_note(slide, "chapter1_analysis_spec_frozen_v1.md; phase1_working_reference.md; ch1_run_config.json")
    add_slide_number(slide, "2")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Cohort Construction And Exclusions", "From standardized ASIC input to the retained Chapter 1 cohort", "Methods")
    add_card(
        slide,
        Inches(0.58),
        Inches(1.18),
        Inches(4.1),
        Inches(1.3),
        "Exact cohort definition",
        [
            "Adult ICU stays with mechanical ventilation >=24h.",
            "First ICU stay via readmission = 0 proxy.",
            "Valid in-ICU mortality label and >=1 valid prediction instance required.",
        ],
    )
    # Flow boxes
    flow_top = Inches(2.75)
    box_w = Inches(2.35)
    box_h = Inches(1.0)
    flow_lefts = [Inches(0.70), Inches(3.32), Inches(5.94), Inches(8.56)]
    flow_titles = [
        f"Input\n{fmt_int(input_hospitals)} hospitals\n{fmt_int(input_stays)} stays",
        f"Site-eligible\n{fmt_int(retained_hospitals)} hospitals\n{fmt_int(stage_site)} stays",
        f"After mech_vent_ge_24h_qc\n{fmt_int(stage_mech)} stays",
        f"Final retained\n{fmt_int(retained_stays)} stays",
    ]
    fills = [PANEL, ACCENT_LIGHT, PANEL, ACCENT_LIGHT]
    for idx, (left, title, fill) in enumerate(zip(flow_lefts, flow_titles, fills)):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, flow_top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = ACCENT_SOFT if fill == ACCENT_LIGHT else LINE
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(title.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.name = BODY_FONT if i > 0 else TITLE_FONT
            run.font.size = Pt(12 if i == 0 else 16)
            run.font.bold = True
            run.font.color.rgb = TEXT if i == 0 else ACCENT
        if idx < len(flow_lefts) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                left + box_w + Inches(0.08),
                flow_top + Inches(0.32),
                Inches(0.34),
                Inches(0.34),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_SOFT
            arrow.line.color.rgb = ACCENT_SOFT
    add_card(
        slide,
        Inches(10.95),
        Inches(1.18),
        Inches(1.82),
        Inches(2.58),
        "Site logic",
        [
            "Include hospital only if ICU mortality is usable.",
            "Require at least 3 of 4 core physiologic groups with dynamic coverage.",
            "Excluded hospitals: asic_UK00, UK01, UK03, UK06.",
        ],
    )
    add_table(slide, Inches(0.82), Inches(4.30), Inches(4.45), Inches(1.75), retained_counts, font_size=10)
    add_card(
        slide,
        Inches(5.55),
        Inches(4.18),
        Inches(3.38),
        Inches(1.95),
        "Major stay-level gates after site restriction",
        [
            "Upstream mech_vent_ge_24h_qc is the main within-site drop mechanism.",
            "Readmission = 0 is the operational first-stay proxy.",
            "No retained stay fails the saved cohort verification checks.",
        ],
    )
    add_card(
        slide,
        Inches(9.12),
        Inches(4.18),
        Inches(3.64),
        Inches(1.95),
        "What to say aloud",
        [
            "Most contraction happens before modeling, at explicit site and stay eligibility steps.",
            "This is a site-restricted ASIC cohort, not every standardized ASIC hospital.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "The retained cohort is shaped primarily by explicit site eligibility and stay-level gating, so any result should be read relative to this restricted analytic population.")
    add_source_note(slide, "chapter1_site_eligibility.csv; chapter1_counts_by_hospital.csv; chapter1_cohort_summary.csv")
    add_slide_number(slide, "3")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Time Representation, Valid-Instance Rule, And Proxy Labels", "Row construction and labelability are part of the analysis design, not just preprocessing detail", "Methods")
    add_card(
        slide,
        Inches(0.58),
        Inches(1.22),
        Inches(4.6),
        Inches(2.3),
        "Valid prediction-instance rule",
        [
            "Completed 8h blocks are the time grid.",
            "Patient must still be alive and in ICU at prediction time.",
            "At least 3 of 4 core physiologic groups must be observed in-block.",
            "Rows remain only when horizon-specific labelability is unambiguous.",
        ],
    )
    add_card(
        slide,
        Inches(0.58),
        Inches(3.72),
        Inches(4.6),
        Inches(1.65),
        "Observed-data gate",
        [
            f"{fmt_int(obs_qc.loc[obs_qc['metric'] == 'block_rows_total', 'value'].iloc[0])} unique usable 8h blocks before horizon duplication.",
            f"{fmt_pct(obs_qc.loc[obs_qc['metric'] == 'value_4_proportion', 'value'].iloc[0], 1)} with all 4 core groups observed.",
            f"{fmt_pct(obs_qc.loc[obs_qc['metric'] == 'value_3_proportion', 'value'].iloc[0], 1)} with exactly 3 of 4.",
        ],
    )
    add_image_contain(slide, assets["slide04_label_diagnostics"], Inches(5.45), Inches(1.18), Inches(7.35), Inches(3.18))
    add_textbox(slide, Inches(5.60), Inches(4.34), Inches(7.0), Inches(0.18), "Figure: labelable/unlabeled burden by horizon plus positive prevalence among labelable rows", font_size=9, color=MUTED)

    proxy_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.42), Inches(4.65), Inches(7.38), Inches(1.48))
    proxy_box.fill.solid()
    proxy_box.fill.fore_color.rgb = PANEL
    proxy_box.line.color.rgb = LINE
    tf = proxy_box.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = [
        "Positive: icu_mortality = 1 and icu_end_time_proxy_hours in (t, t+H].",
        "Negative: icu_mortality = 0 and icu_end_time_proxy_hours >= t+H.",
        "Otherwise unlabeled; 24h counts are 231,596 labelable, 4,986 positive, 226,610 negative, 78,048 unlabeled.",
        "ASIC does not provide true death/discharge timestamps, so these are proxy within-horizon labels.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(11)
            run.font.color.rgb = TEXT
    add_interpretation_box(slide, "Both row retention and labelability are shaped by current-block measurement coverage plus the proxy ICU-end-time rule, so the analyzed row set is intentionally restrictive.", caveat=True)
    add_source_note(slide, "label_logic_audit.md; chapter1_cohort_summary.csv; chapter1_observation_process_qc_summary.csv")
    add_slide_number(slide, "4")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Frozen Modeling Design Choices", "Feature boundary, preprocessing policy, split strategy, baselines, and metrics", "Methods")
    add_card(
        slide,
        Inches(0.58),
        Inches(1.28),
        Inches(3.0),
        Inches(2.05),
        "Feature boundary",
        [
            "Primary feature set: 31 routine variables.",
            "Extended set adds 15 sparse-lab variables.",
            "Model-ready export selects 186 blocked dynamic columns.",
            "98 LOCF/missingness indicator columns are appended.",
        ],
    )
    add_card(
        slide,
        Inches(3.78),
        Inches(1.28),
        Inches(3.0),
        Inches(2.05),
        "Preprocessing policy",
        [
            "Bounded LOCF only for pre-specified families.",
            "Ventilator LOCF only within ventilation-supported windows.",
            "No final imputation in preprocessing export.",
            "Downstream imputation is deferred to model training.",
        ],
    )
    add_card(
        slide,
        Inches(6.98),
        Inches(1.28),
        Inches(2.75),
        Inches(2.05),
        "Split strategy",
        [
            "Target 70 / 15 / 15.",
            "Operational split unit: stay_id_global.",
            "Splitting within retained hospitals.",
            "Frozen seed: 20260327.",
        ],
    )
    add_card(
        slide,
        Inches(9.93),
        Inches(1.28),
        Inches(2.82),
        Inches(2.05),
        "Baselines and metrics",
        [
            "Baselines: logistic regression and XGBoost.",
            "Metrics: AUROC, AUPRC, intercept, slope, Brier.",
            "Calibration is a gating issue before hard-case interpretation.",
        ],
    )
    add_metric_card(
        slide,
        Inches(0.76),
        Inches(3.72),
        Inches(2.15),
        Inches(1.3),
        "Model-ready rows",
        fmt_int(model_ready.loc[model_ready["metric"] == "model_ready_rows_total", "value"].iloc[0]),
        "Across 5 horizons and 3 splits",
    )
    add_metric_card(slide, Inches(3.12), Inches(3.72), Inches(2.15), Inches(1.3), "Base features", "31", "Frozen primary set")
    add_metric_card(slide, Inches(5.48), Inches(3.72), Inches(2.15), Inches(1.3), "Selected columns", "186", "Blocked dynamic feature columns")
    add_metric_card(slide, Inches(7.84), Inches(3.72), Inches(2.15), Inches(1.3), "Indicators", "98", "LOCF and missingness indicators")
    add_metric_card(slide, Inches(10.20), Inches(3.72), Inches(2.15), Inches(1.3), "Export imputation", "None", "Deferred to training stage")
    add_card(
        slide,
        Inches(0.76),
        Inches(5.22),
        Inches(11.6),
        Inches(1.15),
        "Operational caveat",
        [
            "The frozen intent is first-stay analysis, but ASIC lacks patient identifiers, so splitting remains stay-level after first-stay proxy filtering.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "The design intentionally exposes missingness and emphasizes calibration-aware baselines rather than relying on aggressive feature engineering or a large model zoo.")
    add_source_note(slide, "ch1_feature_sets.json; chapter1_primary_readiness_summary.csv; preprocessing_interface.md")
    add_slide_number(slide, "5")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Retained Cohort And Realized Split Summary", "Current ASIC run scale and the realized internal split balance", "Results")
    add_metric_card(slide, Inches(0.60), Inches(1.18), Inches(2.35), Inches(1.38), "Retained stays", fmt_int(retained_stays), "Final Chapter 1 cohort")
    add_metric_card(slide, Inches(3.10), Inches(1.18), Inches(2.35), Inches(1.38), "Valid instances", fmt_int(valid_total), "Across all five horizons")
    add_metric_card(slide, Inches(5.60), Inches(1.18), Inches(2.35), Inches(1.38), "Per horizon", fmt_int(valid_per_h), "Before labelability filtering")
    add_metric_card(slide, Inches(8.10), Inches(1.18), Inches(2.35), Inches(1.38), "24h test rows", fmt_int(33676), "With 752 positives")
    add_metric_card(slide, Inches(10.60), Inches(1.18), Inches(2.15), Inches(1.38), "Test stays", fmt_int(967), "26.78% stay mortality")
    add_textbox(slide, Inches(0.68), Inches(2.85), Inches(4.1), Inches(0.22), "Realized stay split", font_size=14, bold=True)
    add_table(slide, Inches(0.64), Inches(3.10), Inches(5.45), Inches(1.8), split_table, font_size=10)
    add_textbox(slide, Inches(6.35), Inches(2.85), Inches(4.5), Inches(0.22), "24h row-level split summary", font_size=14, bold=True)
    add_table(slide, Inches(6.28), Inches(3.10), Inches(6.0), Inches(1.8), split_24, font_size=10)
    add_card(
        slide,
        Inches(0.80),
        Inches(5.18),
        Inches(11.35),
        Inches(1.2),
        "What this changes relative to the synthetic local bundle",
        [
            "The current ASIC package is a full-scale internal run with a binary-evaluable test split.",
            "Main-deck performance numbers can therefore be reported on test rather than falling back to validation.",
        ],
    )
    add_interpretation_box(slide, "The current ASIC bundle is a real internal evaluation package, not a local smoke-test artifact, which materially changes how performance slides should be read.")
    add_source_note(slide, "chapter1_cohort_summary.csv; chapter1_stay_split_summary.csv; chapter1_primary_split_summary.csv")
    add_slide_number(slide, "6")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Baseline Model Performance And Calibration", "Headline quantitative baseline results on the true test split", "Results")
    add_image_contain(slide, assets["slide07_baseline_overview"], Inches(0.62), Inches(1.22), Inches(6.5), Inches(4.55))
    add_textbox(slide, Inches(0.72), Inches(5.72), Inches(6.0), Inches(0.2), "Figure: test AUROC and calibration slope across horizons for logistic regression and XGBoost", font_size=9, color=MUTED)
    add_textbox(slide, Inches(7.35), Inches(1.22), Inches(4.7), Inches(0.22), "24h test metrics", font_size=15, bold=True)
    add_table(slide, Inches(7.25), Inches(1.52), Inches(5.35), Inches(1.82), metrics_24_table, font_size=10)
    add_card(
        slide,
        Inches(7.28),
        Inches(3.60),
        Inches(5.30),
        Inches(1.6),
        "Readout",
        [
            "Logistic regression: AUROC 0.819, AUPRC 0.268, slope 0.974, Brier 0.0186.",
            "XGBoost: AUROC 0.848, AUPRC 0.318, slope 1.162, Brier 0.1351.",
            "XGBoost ranks better; logistic regression remains the cleaner probability model.",
        ],
    )
    add_interpretation_box(slide, "The ranking-versus-calibration tradeoff is real in this ASIC run: XGBoost discriminates better, while logistic regression is the more stable anchor for calibration-first hard-case interpretation.")
    add_source_note(slide, "combined_metrics.csv; reporting_split_summary.csv")
    add_slide_number(slide, "7")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Primary 24h Mortality-Vs-Risk Structure", "Logistic regression at the frozen primary horizon", "Results")
    add_image_contain(
        slide,
        CLUSTER_DIR
        / "evaluation/asic/baselines/primary_medians/logistic_regression/horizon_24h/mortality_vs_risk_plot.png",
        Inches(0.58),
        Inches(1.18),
        Inches(7.55),
        Inches(5.18),
    )
    add_image_contain(
        slide,
        CLUSTER_DIR
        / "evaluation/asic/baselines/primary_medians/logistic_regression/horizon_24h/reliability_plot.png",
        Inches(8.32),
        Inches(1.18),
        Inches(4.42),
        Inches(5.18),
    )
    add_textbox(slide, Inches(0.66), Inches(6.30), Inches(7.2), Inches(0.2), "Figure: 24h observed mortality across risk bins for logistic regression on the test split", font_size=9, color=MUTED)
    add_textbox(slide, Inches(8.42), Inches(6.30), Inches(4.0), Inches(0.2), "Figure: 24h reliability plot", font_size=9, color=MUTED)
    add_interpretation_box(slide, "At 24h the logistic model preserves coherent risk ordering with acceptable calibration, which makes the hard-case analysis interpretable as a bounded descriptive exercise.")
    add_source_note(slide, "horizon_24h/mortality_vs_risk_plot.png; horizon_24h/reliability_plot.png")
    add_slide_number(slide, "8")

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Hard-Case Definition And Burden", "Operational rule and resulting low-predicted fatal burden by horizon", "Results")
    add_card(
        slide,
        Inches(0.62),
        Inches(1.18),
        Inches(4.25),
        Inches(1.8),
        "Frozen hard-case rule",
        [
            "Rule: asic_logistic_last_eligible_nonfatal_q75_v1.",
            "Collapse to one last eligible stay-level point per stay and horizon.",
            "Mark fatal stays as hard cases if predicted risk is below the nonfatal q75 threshold.",
        ],
    )
    add_metric_card(slide, Inches(0.78), Inches(3.28), Inches(1.9), Inches(1.25), "24h fatal stays", fmt_int(1682), "Last eligible stay-level slice")
    add_metric_card(slide, Inches(2.88), Inches(3.28), Inches(1.9), Inches(1.25), "24h hard cases", fmt_int(346), "Low-predicted fatal stays")
    add_image_contain(slide, assets["slide09_hard_case_burden"], Inches(5.15), Inches(1.18), Inches(7.45), Inches(4.25))
    add_textbox(slide, Inches(5.30), Inches(5.42), Inches(6.5), Inches(0.2), "Figure: hard-case share among fatal stays across the frozen horizon grid", font_size=9, color=MUTED)
    add_table(slide, Inches(0.78), Inches(4.82), Inches(4.0), Inches(1.25), thresholds, font_size=10)
    add_interpretation_box(slide, "Under the frozen logistic rule, low-predicted fatal stays are a substantial minority of fatal stays at every horizon rather than a rare edge case.")
    add_source_note(slide, "horizon_hard_case_summary.csv")
    add_slide_number(slide, "9")

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_header(slide, "24h Hard-Case Comparison", "Low-predicted fatal stays versus other fatal stays at the primary horizon", "Results")
    add_image_contain(
        slide,
        CLUSTER_DIR
        / "evaluation/asic/hard_cases/primary_medians/logistic_regression/asic_hard_case_comparison/effect_size_figure.png",
        Inches(0.55),
        Inches(1.16),
        Inches(8.65),
        Inches(5.55),
    )
    add_card(
        slide,
        Inches(9.45),
        Inches(1.22),
        Inches(3.15),
        Inches(1.35),
        "24h exported summary slice",
        [
            "Fatal stays: 1,682.",
            "Low-predicted fatal: 346.",
            "Other fatal: 1,336.",
        ],
    )
    add_card(
        slide,
        Inches(9.45),
        Inches(2.72),
        Inches(3.15),
        Inches(2.02),
        "Strongest differences",
        [
            f"Prediction time: {pred_time_row['low_predicted_fatal_stays']} vs {pred_time_row['other_fatal_stays']}.",
            f"MAP: {map_row['low_predicted_fatal_stays']} vs {map_row['other_fatal_stays']}.",
            f"PF ratio: {pf_row['low_predicted_fatal_stays']} vs {pf_row['other_fatal_stays']}.",
            f"PEEP: {peep_row['low_predicted_fatal_stays']} vs {peep_row['other_fatal_stays']}.",
        ],
    )
    add_card(
        slide,
        Inches(9.45),
        Inches(4.95),
        Inches(3.15),
        Inches(1.35),
        "Bounded subgroup signals",
        [
            f"Hospital: {hospital_row['low_predicted_fatal_stays']}.",
            f"Disease group: {disease_row['low_predicted_fatal_stays']}.",
            "Exact age is unavailable; age_group only.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(
        slide,
        "Low-predicted fatal stays look less aligned with captured short-term physiologic severity, but this remains a descriptive operational comparison rather than a typology of death. The local presentation uses the approved aggregate export bundle, not the restricted row-level reconstruction table.",
    )
    add_source_note(
        slide,
        "approved aggregate export bundle: effect_size_figure.png; comparison_table.csv; summary.md; asic_hard_case_comparison_variable_audit_memo.md",
    )
    add_slide_number(slide, "10")

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Horizon Dependence", "Does the hard-case burden survive across the frozen horizon grid?", "Results")
    add_image_contain(
        slide,
        CLUSTER_DIR / "evaluation/asic/horizon_dependence/final/mortality_risk_horizon_comparison.png",
        Inches(0.52),
        Inches(1.18),
        Inches(12.2),
        Inches(4.6),
    )
    add_textbox(slide, Inches(0.70), Inches(5.74), Inches(7.0), Inches(0.2), "Figure: mortality-versus-risk structure across 8h, 16h, 24h, 48h, and 72h horizons", font_size=9, color=MUTED)
    add_metric_card(slide, Inches(0.72), Inches(5.96), Inches(2.45), Inches(0.95), "Mean pairwise Jaccard", f"{overlap_mean:.3f}", "Across all horizon pairs")
    add_metric_card(slide, Inches(3.38), Inches(5.96), Inches(2.45), Inches(0.95), "24h vs 48h", f"{overlap_24_48:.3f}", "Hard-case overlap")
    add_metric_card(slide, Inches(6.04), Inches(5.96), Inches(2.45), Inches(0.95), "24h vs 72h", f"{overlap_24_72:.3f}", "Hard-case overlap")
    add_metric_card(slide, Inches(8.70), Inches(5.96), Inches(1.75), Inches(0.95), "Manifest label", horizon_manifest["interpretation_label"], "Current saved read")
    add_card(
        slide,
        Inches(10.65),
        Inches(5.90),
        Inches(2.05),
        Inches(1.0),
        "Note",
        ["The figure is descriptive.", "It does not redefine the hard-case rule."],
    )
    add_interpretation_box(slide, "The current ASIC outputs support a descriptive persistence read: the hard-case burden changes little across horizons and cross-horizon membership overlap remains high.")
    add_source_note(slide, "mortality_risk_horizon_comparison.png; pairwise_overlap.csv; final/run_manifest.json")
    add_slide_number(slide, "11")

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Secondary Robustness Summary", "Cross-model agreement is mature; temporal aggregation remains provisional", "Results")
    add_card(
        slide,
        Inches(0.60),
        Inches(1.18),
        Inches(4.25),
        Inches(1.75),
        "Cross-model agreement at 24h",
        [
            "Logistic hard: 346.",
            "XGBoost-Platt hard: 227.",
            "Overlap: 188.",
            "Jaccard: 0.488.",
            "54% of logistic hard cases are confirmed by XGBoost-Platt.",
        ],
    )
    add_card(
        slide,
        Inches(0.60),
        Inches(3.10),
        Inches(4.25),
        Inches(1.85),
        "Temporal preview: 8h vs 16h aggregation",
        [
            "Logistic 24h AUROC 0.819 -> 0.816; AUPRC 0.268 -> 0.235.",
            "XGBoost 24h AUROC 0.848 -> 0.846; AUPRC 0.318 -> 0.291.",
            "Encouraging, but still only a preview rather than a formal sensitivity run.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_card(
        slide,
        Inches(0.60),
        Inches(5.16),
        Inches(4.25),
        Inches(1.18),
        "Site sanity check",
        [
            "Pooled 24h result is not obviously driven by one hospital.",
            "Site-level calibration is still much sparser and more variable than pooled results.",
        ],
    )
    add_image_contain(
        slide,
        CLUSTER_DIR / "evaluation/asic/baselines/primary_medians/logistic_regression/primary_24h_site_overview.png",
        Inches(5.10),
        Inches(1.10),
        Inches(7.55),
        Inches(5.35),
    )
    add_textbox(slide, Inches(5.30), Inches(6.35), Inches(6.5), Inches(0.2), "Figure: 24h site-stratified logistic overview for the test split", font_size=9, color=MUTED)
    add_interpretation_box(slide, "Robustness is mixed rather than all-or-none: horizon persistence is strong, cross-model agreement is only moderate, and temporal aggregation has an encouraging but incomplete preview.")
    add_source_note(slide, "horizon_hard_case_agreement_summary.csv; primary_24h_site_overview.png; aggregation_comparison_metrics.csv")
    add_slide_number(slide, "12")

    # Slide 13
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Current Bounded Interpretation", "Most defensible technical readout of the current ASIC package", "Results")
    add_card(
        slide,
        Inches(0.72),
        Inches(1.26),
        Inches(5.85),
        Inches(4.9),
        "Supported now",
        [
            "ASIC supports heterogeneous predictability under the observed feature set.",
            "Logistic 24h risk is a usable calibration-aware anchor for hard-case definition.",
            "About one-fifth of fatal stays are low-predicted under the frozen logistic rule.",
            "The low-predicted fatal burden persists descriptively across the frozen horizons.",
            "Hard cases are operational and measurement-bound, not ontological classes.",
        ],
    )
    add_card(
        slide,
        Inches(6.82),
        Inches(1.26),
        Inches(5.78),
        Inches(4.9),
        "Explicit non-claims",
        [
            "No biological death subtype claim.",
            "No claim that these deaths are inherently or irreducibly unpredictable.",
            "No causal explanation for why cases are low-predicted.",
            "No claim that hard-case membership is model-invariant at patient level.",
            "No external-transport claim, because MIMIC results are not available yet.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "The defensible claim is about heterogeneous predictability under this recorded representation; the deck should not slide into subtype language or deeper causal interpretation.")
    add_source_note(slide, "context_sprint4.md; chapter1_analysis_spec_frozen_v1.md; hard-case comparison summary.md")
    add_slide_number(slide, "13")

    # Slide 14
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Open Limitations And Pending Analyses", "What still blocks a more closed ASIC interpretation", "Results")
    add_card(
        slide,
        Inches(0.62),
        Inches(1.20),
        Inches(5.85),
        Inches(4.95),
        "Current limitations",
        [
            "ASIC lacks true death and ICU-discharge timestamps for within-horizon labels.",
            "Rows are shaped by valid-instance and labelability gates, not just model fitting.",
            "Splitting remains stay-level because ASIC lacks patient identifiers.",
            "The 24h hard-case package lacks exact age; age_group only is available.",
            "Several comparison proxies depend on partial coverage or bounded LOCF.",
            "Some older markdown interpretation memos still contain stale preliminary wording.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_card(
        slide,
        Inches(6.72),
        Inches(1.20),
        Inches(5.95),
        Inches(4.95),
        "Pending sensitivity work",
        [
            "Observation-process hard-case comparison.",
            "Treatment-limitation sensitivity or a formal explicit absence note.",
            "Formal temporal-aggregation sensitivity beyond the current preview.",
            "Disease-stratified predictability analyses.",
            "MIMIC external validation and transportability check.",
        ],
    )
    add_interpretation_box(slide, "The current ASIC package is strong enough for a bounded methods/results presentation, but it is not yet closed on the key interpretation-critical sensitivities.", caveat=True)
    add_source_note(slide, "label_logic_audit.md; asic_hard_case_comparison_variable_audit_memo.md; preview_note.md")
    add_slide_number(slide, "14")

    # Appendix A1
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A1 — Full Metrics By Model And Horizon", "Test-split metrics only", "Appendix")
    add_textbox(slide, Inches(0.62), Inches(1.08), Inches(5.0), Inches(0.22), "Logistic regression", font_size=14, bold=True)
    add_table(slide, Inches(0.58), Inches(1.34), Inches(5.85), Inches(4.85), full_metrics_log, font_size=9)
    add_textbox(slide, Inches(6.88), Inches(1.08), Inches(5.0), Inches(0.22), "XGBoost", font_size=14, bold=True)
    add_table(slide, Inches(6.82), Inches(1.34), Inches(5.95), Inches(4.85), full_metrics_xgb, font_size=9)
    add_interpretation_box(slide, "The ranking-versus-calibration tradeoff is consistent across horizons rather than being unique to 24h.")
    add_source_note(slide, "combined_metrics.csv")
    add_slide_number(slide, "A1")

    # Appendix A2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A2 — Site-Stratified 24h Sanity Check", "Backup view for pooled 24h interpretation", "Appendix")
    add_image_contain(
        slide,
        CLUSTER_DIR / "evaluation/asic/baselines/primary_medians/logistic_regression/primary_24h_site_overview.png",
        Inches(0.55),
        Inches(1.15),
        Inches(7.5),
        Inches(5.3),
    )
    add_textbox(slide, Inches(0.72), Inches(6.35), Inches(6.5), Inches(0.2), "Figure: 24h site overview for logistic regression on the test split", font_size=9, color=MUTED)
    add_table(slide, Inches(8.35), Inches(1.45), Inches(4.2), Inches(2.35), site_24_logistic, font_size=10)
    add_card(
        slide,
        Inches(8.35),
        Inches(4.05),
        Inches(4.2),
        Inches(1.5),
        "Backup read",
        [
            "All four retained hospitals have evaluable 24h test metrics.",
            "Site-level calibration varies more than pooled calibration and should stay secondary.",
        ],
    )
    add_interpretation_box(slide, "The pooled 24h signal is not obviously driven by a single hospital, but site-level inference should remain cautious because counts are much smaller than in the pooled analysis.")
    add_source_note(slide, "primary_24h_site_overview.png; combined_primary_site_summary.csv")
    add_slide_number(slide, "A2")

    # Appendix A3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A3 — Cross-Model Hard-Case Agreement By Horizon", "Logistic regression versus XGBoost-Platt", "Appendix")
    add_table(slide, Inches(0.76), Inches(1.50), Inches(7.35), Inches(2.5), agreement_table, font_size=10)
    add_card(
        slide,
        Inches(8.55),
        Inches(1.48),
        Inches(3.75),
        Inches(1.55),
        "24h anchor",
        [
            "Jaccard 0.488 at 24h.",
            "188 stays are hard in both models.",
            "Broader burden is more stable than exact patient-level membership.",
        ],
    )
    add_card(
        slide,
        Inches(8.55),
        Inches(3.26),
        Inches(3.75),
        Inches(1.85),
        "Pattern across horizons",
        [
            "Agreement improves somewhat at longer horizons.",
            "It remains well short of model invariance.",
            "This weakens decomposition-style claims more than the descriptive hard-case burden itself.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "Cross-model overlap is meaningful but incomplete, so patient-level hard-case membership should be treated as model-sensitive.")
    add_source_note(slide, "horizon_hard_case_agreement_summary.csv")
    add_slide_number(slide, "A3")

    # Appendix A4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A4 — Variable Audit For The 24h Hard-Case Comparison", "Availability and proxy-readiness check", "Appendix")
    add_table(slide, Inches(0.58), Inches(1.32), Inches(12.0), Inches(4.55), variable_audit_table, font_size=9)
    add_card(
        slide,
        Inches(0.75),
        Inches(6.00),
        Inches(11.8),
        Inches(0.78),
        "Bottom line",
        [
            "The comparison variable package is usable, but it is not a complete clinical adjustment set and exact age remains unavailable in the current ASIC layer.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_source_note(slide, "asic_hard_case_comparison_variable_audit_memo.md")
    add_slide_number(slide, "A4")

    # Appendix A5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A5 — Why SOFA Was Not Used", "SOFA feasibility audit for the 24h fatal-stay slice", "Appendix")
    add_card(slide, Inches(0.62), Inches(1.25), Inches(4.0), Inches(2.2), "Missing standard components", ["No GCS / CNS field.", "No vasopressor variables in the model-ready layer.", "No urine-output field for renal SOFA.", "Standard SOFA scoring therefore cannot be completed."])
    add_card(slide, Inches(4.82), Inches(1.25), Inches(3.9), Inches(2.2), "Available but incomplete", ["Respiratory scorable in 73%.", "Platelets 92%, but heavily LOCF-based.", "Bilirubin 71%, with major 48h LOCF dependence.", "Creatinine 98%, but 65% depends on 48h LOCF."])
    add_card(slide, Inches(8.92), Inches(1.25), Inches(3.7), Inches(2.2), "Complete-case problem", ["Only 244/1,682 (15%) complete cases with current-block values.", "825/1,682 (49%) even after LOCF.", "Coverage is structured, not plausibly MCAR."])
    add_bullets(slide, Inches(0.86), Inches(4.05), Inches(11.0), Inches(1.6), sofa_points, font_size=13)
    add_interpretation_box(slide, "A direct proxy table is cleaner and more reproducible than an incomplete pseudo-SOFA built from missing standard domains.")
    add_source_note(slide, "sofa_feasibility_memo.md")
    add_slide_number(slide, "A5")

    # Appendix A6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A6 — Observation-Process Readiness", "Instrumentation exists; the explanatory sensitivity analysis does not yet", "Appendix")
    add_card(
        slide,
        Inches(0.62),
        Inches(1.22),
        Inches(5.1),
        Inches(2.65),
        "Derived variables",
        [
            "obs_hr_grp_block, obs_bp_grp_block, obs_resp_grp_block, obs_oxy_grp_block.",
            "n_core_grps_obs_block.",
            "tsl_hr_grp_h, tsl_bp_grp_h, tsl_resp_grp_h, tsl_oxy_grp_h.",
            "Block-level export is unique before horizon duplication.",
        ],
    )
    add_table(slide, Inches(6.05), Inches(1.35), Inches(6.2), Inches(2.8), obs_summary_table, font_size=10)
    add_card(
        slide,
        Inches(0.74),
        Inches(4.25),
        Inches(11.5),
        Inches(1.55),
        "Current status",
        [
            "Observation-process variables are derived and QC'd, but no saved hard-case comparison using them is present in the current result bundle.",
            "That means the instrumentation is ready while the interpretation-critical sensitivity result remains pending.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "Observation-process confounding remains an open interpretive risk because readiness is documented but the actual sensitivity result is still missing.")
    add_source_note(slide, "chapter1_observation_process_qc_summary.csv; chapter1_observation_process_implementation_note.md")
    add_slide_number(slide, "A6")

    # Appendix A7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix A7 — Temporal Aggregation Preview: 8h Vs 16h", "Provisional backup slide, not a formal sensitivity conclusion", "Appendix")
    add_image_contain(
        slide,
        CLUSTER_DIR / "temporal_preview/asic/aggregation_16h/comparison/logistic_regression_24h_mortality_vs_risk_8h_vs_16h.png",
        Inches(0.55),
        Inches(1.15),
        Inches(7.75),
        Inches(4.95),
    )
    add_textbox(slide, Inches(0.72), Inches(6.00), Inches(7.0), Inches(0.2), "Figure: logistic 24h mortality-vs-risk structure under 8h versus 16h aggregation", font_size=9, color=MUTED)
    add_table(slide, Inches(8.55), Inches(1.45), Inches(4.0), Inches(1.7), temporal_table, font_size=10)
    add_card(
        slide,
        Inches(8.55),
        Inches(3.42),
        Inches(4.0),
        Inches(1.85),
        "Preview note",
        [
            "Maximum absolute AUROC delta across comparable holdout pairs: 0.018.",
            "Maximum absolute AUPRC delta: 0.066.",
            "The preview does not choose an optimal aggregation and should not refreeze the design alone.",
        ],
        fill=RUST_LIGHT,
        title_fill=RGBColor(254, 215, 170),
    )
    add_interpretation_box(slide, "The early temporal preview shows movement without obvious collapse, which is encouraging but still too narrow to close the temporal-sensitivity question.")
    add_source_note(slide, "logistic_regression_24h_mortality_vs_risk_8h_vs_16h.png; aggregation_comparison_metrics.csv; preview_note.md")
    add_slide_number(slide, "A7")

    prs.save(OUTPUT_PPTX)


def main() -> None:
    ensure_dirs()
    data = load_data()
    assets = {
        "slide04_label_diagnostics": export_label_diagnostics(data["cohort_summary"]),
        "slide07_baseline_overview": export_baseline_overview(data["combined_metrics"]),
        "slide09_hard_case_burden": export_hard_case_burden(data["hard_case_summary"]),
    }
    build_presentation(data, assets)
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
